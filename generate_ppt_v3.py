#!/usr/bin/env python3
"""Generate an editable PowerPoint slide for AVD Private Link architecture (v3).
- Compact, smaller icons to avoid overlap
- Properly calculated arrow connections
- Corrected auth flow: Client → Entra ID (public) → then SD-WAN (private)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── Colors ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x66, 0x66, 0x66)
DARK_BLUE = RGBColor(0x0D, 0x47, 0xA1)
BLUE = RGBColor(0x15, 0x65, 0xC0)
LIGHT_BLUE = RGBColor(0x1E, 0x88, 0xE5)
SKY_BLUE = RGBColor(0x42, 0xA5, 0xF5)
OCEAN_BLUE = RGBColor(0x02, 0x77, 0xBD)
DEEP_OCEAN = RGBColor(0x01, 0x57, 0x9B)
TEAL = RGBColor(0x03, 0x9B, 0xE5)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
DARK_ORANGE = RGBColor(0xBF, 0x36, 0x0C)
LIGHT_ORANGE = RGBColor(0xEF, 0x6C, 0x00)
WARM_ORANGE = RGBColor(0xF4, 0x7B, 0x13)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
DARK_GREEN = RGBColor(0x1B, 0x5E, 0x20)
MID_GREEN = RGBColor(0x38, 0x8E, 0x3C)
FOREST = RGBColor(0x00, 0x60, 0x00)
YELLOW = RGBColor(0xF9, 0xA8, 0x25)
AMBER = RGBColor(0xE6, 0x51, 0x00)
STEP_BLUE = RGBColor(0x1A, 0x73, 0xE8)
STEP_RED = RGBColor(0xD5, 0x00, 0x00)
SHORTPATH_ORANGE = RGBColor(0xFF, 0x6D, 0x00)
SDWAN_GREEN = RGBColor(0x00, 0xC8, 0x53)
ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def rect(l, t, w, h, fill, border, text="", fs=9, fc=WHITE, bold=False, bw=Pt(2)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = bw
    tf = s.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(fs)
    r.font.color.rgb = fc; r.font.bold = bold
    return s


def txt(l, t, w, h, text, fs=9, fc=BLACK, bold=False, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(fs)
    r.font.color.rgb = fc; r.font.bold = bold
    return tb


def oval(l, t, w, h, fill, border, text="", fs=8, fc=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(2)
    tf = s.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(fs)
    r.font.color.rgb = fc; r.font.bold = True
    return s


def arrow(x1, y1, x2, y2, color, w=Pt(2), dashed=False):
    c = slide.shapes.add_connector(1, x1, y1, x2, y2)
    c.line.color.rgb = color; c.line.width = w
    c.line.fill.solid(); c.line.color.rgb = color
    ln = c.line._get_or_add_ln()
    te = etree.SubElement(ln, f'{{{ns}}}tailEnd')
    te.set('type', 'triangle'); te.set('w', 'med'); te.set('len', 'med')
    if dashed:
        pd = etree.SubElement(ln, f'{{{ns}}}prstDash'); pd.set('val', 'dash')
    return c


def step(x, y, num, label, color=STEP_BLUE):
    oval(x, y, Inches(0.28), Inches(0.28), color, color, num, fs=8, fc=WHITE)
    txt(x + Inches(0.32), y - Inches(0.01), Inches(1.4), Inches(0.30),
        label, fs=7, fc=BLACK, bold=False)


# ══════════════════════════════════════════════════════════
# TITLE
# ══════════════════════════════════════════════════════════
txt(Inches(0.3), Inches(0.1), Inches(12), Inches(0.4),
    "Azure Virtual Desktop — Private Link Architecture (v3)",
    fs=20, fc=DARK_BLUE, bold=True, align=PP_ALIGN.LEFT)
txt(Inches(0.3), Inches(0.45), Inches(10), Inches(0.25),
    "End-to-end private connectivity via SD-WAN + AVD Private Link  |  Azure China (21Vianet)",
    fs=9, fc=GRAY, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════
# LAYOUT GRID (all coordinates calculated from here)
# ══════════════════════════════════════════════════════════
# Slide: 13.333 x 7.5
# Left: OnPrem (0.2 - 2.1)  |  SD-WAN circle (2.8)  |  CE2 (4.0 - 7.8)  |  CN3 (8.3 - 12.1)
# Top margin: 0.85

MARGIN_T = Inches(0.85)
REGION_H = Inches(5.5)

# ── On-Premises ──
OP_L, OP_W = Inches(0.2), Inches(1.9)
OP_T = MARGIN_T
rect(OP_L, OP_T, OP_W, REGION_H, GREEN, DARK_GREEN, bw=Pt(3))
txt(OP_L, OP_T + Inches(0.05), OP_W, Inches(0.3),
    "🏢 On-Premises", fs=9, fc=WHITE, bold=True)

# Shenzhen
SZ_L, SZ_T = OP_L + Inches(0.15), OP_T + Inches(0.45)
SZ_W, SZ_H = Inches(1.6), Inches(0.9)
rect(SZ_L, SZ_T, SZ_W, SZ_H, MID_GREEN, DARK_GREEN,
     "🖥️ Thin Clients\n(Shenzhen)\n10.1.1.0/24", fs=7, bw=Pt(1))

# Dalian
DL_L, DL_T = OP_L + Inches(0.15), OP_T + Inches(1.55)
DL_W, DL_H = Inches(1.6), Inches(0.9)
rect(DL_L, DL_T, DL_W, DL_H, MID_GREEN, DARK_GREEN,
     "🖥️ Thin Clients\n(Dalian)\n10.2.1.0/24", fs=7, bw=Pt(1))

# SD-WAN placeholder text
SDWAN_CX, SDWAN_CY = Inches(2.85), OP_T + Inches(3.2)
SDWAN_R = Inches(0.55)
oval(SDWAN_CX - SDWAN_R, SDWAN_CY - SDWAN_R, SDWAN_R * 2, SDWAN_R * 2,
     FOREST, DARK_GREEN, "🔒\nSD-WAN", fs=7, fc=WHITE)

# ── Azure China East 2 ──
CE2_L, CE2_W = Inches(4.0), Inches(3.8)
CE2_T = MARGIN_T
rect(CE2_L, CE2_T, CE2_W, REGION_H, BLUE, DARK_BLUE, bw=Pt(3))
txt(CE2_L, CE2_T + Inches(0.05), CE2_W, Inches(0.3),
    "☁️ Azure China East 2", fs=10, fc=WHITE, bold=True)

# VNet Gateway
GW_L, GW_T = CE2_L + Inches(0.3), CE2_T + Inches(0.4)
GW_W, GW_H = Inches(3.2), Inches(0.45)
rect(GW_L, GW_T, GW_W, GW_H, DARK_BLUE, DEEP_OCEAN,
     "🔗 VNet Gateway · SD-WAN Hub", fs=9, fc=WHITE, bold=True, bw=Pt(2))

# AVD Control Plane sub-box
AVD_L, AVD_T = CE2_L + Inches(0.2), CE2_T + Inches(1.1)
AVD_W, AVD_H = Inches(3.4), Inches(1.7)
rect(AVD_L, AVD_T, AVD_W, AVD_H, LIGHT_BLUE, BLUE, bw=Pt(1))
txt(AVD_L, AVD_T + Inches(0.03), AVD_W, Inches(0.25),
    "AVD Control Plane · Private Link", fs=8, fc=WHITE, bold=True)

# PE Workspace
PEWS_L, PEWS_T = AVD_L + Inches(0.1), AVD_T + Inches(0.35)
PEWS_W, PEWS_H = Inches(1.5), Inches(1.15)
rect(PEWS_L, PEWS_T, PEWS_W, PEWS_H, SKY_BLUE, LIGHT_BLUE,
     "🔗 Private Endpoint\nWorkspace\nFeed Discovery", fs=7, bw=Pt(1))

# PE Host Pool
PEHP_L, PEHP_T = AVD_L + Inches(1.75), AVD_T + Inches(0.35)
PEHP_W, PEHP_H = Inches(1.5), Inches(1.15)
rect(PEHP_L, PEHP_T, PEHP_W, PEHP_H, SKY_BLUE, LIGHT_BLUE,
     "🔗 Private Endpoint\nHost Pool\nRDP Connection", fs=7, bw=Pt(1))

# Business App VNet
BIZ_L, BIZ_T = CE2_L + Inches(0.2), CE2_T + Inches(3.05)
BIZ_W, BIZ_H = Inches(3.4), Inches(2.2)
rect(BIZ_L, BIZ_T, BIZ_W, BIZ_H, OCEAN_BLUE, DEEP_OCEAN, bw=Pt(1))
txt(BIZ_L, BIZ_T + Inches(0.03), BIZ_W, Inches(0.25),
    "Business App VNet · 10.10.0.0/16", fs=8, fc=WHITE, bold=True)

# App Server
APP_L, APP_T = BIZ_L + Inches(0.3), BIZ_T + Inches(0.4)
APP_W, APP_H = Inches(2.8), Inches(0.7)
rect(APP_L, APP_T, APP_W, APP_H, TEAL, OCEAN_BLUE,
     "🖧 Business Application Server", fs=9, bw=Pt(1))

# VNet Peering label in CE2
rect(BIZ_L + Inches(0.3), BIZ_T + Inches(1.3), Inches(2.8), Inches(0.65),
     DEEP_OCEAN, OCEAN_BLUE,
     "🔗 VNet Peering ↔ China North 3", fs=8, bw=Pt(1))

# ── Azure China North 3 ──
CN3_L, CN3_W = Inches(8.3), Inches(3.8)
CN3_T = MARGIN_T
rect(CN3_L, CN3_T, CN3_W, REGION_H, ORANGE, DARK_ORANGE, bw=Pt(3))
txt(CN3_L, CN3_T + Inches(0.05), CN3_W, Inches(0.3),
    "☁️ Azure China North 3", fs=10, fc=WHITE, bold=True)

# Session Host VNet
SHV_L, SHV_T = CN3_L + Inches(0.15), CN3_T + Inches(0.4)
SHV_W, SHV_H = Inches(3.5), Inches(4.85)
rect(SHV_L, SHV_T, SHV_W, SHV_H, LIGHT_ORANGE, ORANGE, bw=Pt(1))
txt(SHV_L, SHV_T + Inches(0.03), SHV_W, Inches(0.25),
    "Session Host VNet · 10.20.0.0/16", fs=8, fc=WHITE, bold=True)

# Session Host VM
SH_L, SH_T = SHV_L + Inches(0.15), SHV_T + Inches(0.35)
SH_W, SH_H = Inches(3.2), Inches(0.9)
rect(SH_L, SH_T, SH_W, SH_H, WARM_ORANGE, ORANGE,
     "🖥️ Session Host VM\nDesktop & RemoteApp", fs=9, fc=WHITE, bold=True, bw=Pt(1))

# VNet Peering (CN3)
rect(SHV_L + Inches(0.15), SHV_T + Inches(1.45), Inches(3.2), Inches(0.6),
     LIGHT_ORANGE, ORANGE, "🔗 VNet Peering ↔ China East 2", fs=8, bw=Pt(1))

# Private DNS Zone
rect(SHV_L + Inches(0.15), SHV_T + Inches(2.25), Inches(3.2), Inches(0.6),
     DARK_ORANGE, ORANGE, "🔒 Private DNS Zone\nprivatelink.wvd.azure.cn", fs=8, bw=Pt(1))

# NSG
rect(SHV_L + Inches(0.15), SHV_T + Inches(3.05), Inches(3.2), Inches(0.6),
     DARK_ORANGE, ORANGE, "🛡️ NSG · Outbound Rules\nAllow UDP 3390", fs=8, bw=Pt(1))

# RDP Shortpath label
rect(SHV_L + Inches(0.15), SHV_T + Inches(3.85), Inches(3.2), Inches(0.6),
     WARM_ORANGE, DARK_ORANGE, "⚡ RDP Shortpath · UDP 3390\nManaged Networks", fs=8, bw=Pt(1))

# ── Entra ID (top center) ──
ENTRA_L, ENTRA_T = Inches(1.0), Inches(0.08)
ENTRA_W, ENTRA_H = Inches(2.2), Inches(0.5)
rect(ENTRA_L, ENTRA_T, ENTRA_W, ENTRA_H, YELLOW, AMBER,
     "🌐 Entra ID (Public)", fs=8, fc=BLACK, bold=True, bw=Pt(2))

# ══════════════════════════════════════════════════════════
# ARROWS — all calculated from shape positions
# ══════════════════════════════════════════════════════════

# Centers / edges
sz_cx, sz_cy = SZ_L + SZ_W / 2, SZ_T + SZ_H / 2
dl_cx, dl_cy = DL_L + DL_W / 2, DL_T + DL_H / 2
sz_right = SZ_L + SZ_W
dl_right = DL_L + DL_W
op_right = OP_L + OP_W

sdwan_left = SDWAN_CX - SDWAN_R
sdwan_right = SDWAN_CX + SDWAN_R

gw_cx = GW_L + GW_W / 2
gw_cy = GW_T + GW_H / 2
gw_left = GW_L
gw_bottom = GW_T + GW_H
gw_right = GW_L + GW_W

pews_cx = PEWS_L + PEWS_W / 2
pews_cy = PEWS_T + PEWS_H / 2
pews_top = PEWS_T
pehp_cx = PEHP_L + PEHP_W / 2
pehp_cy = PEHP_T + PEHP_H / 2
pehp_right = PEHP_L + PEHP_W

sh_left = SH_L
sh_cx = SH_L + SH_W / 2
sh_cy = SH_T + SH_H / 2
sh_top = SH_T

app_cx = APP_L + APP_W / 2
app_cy = APP_T + APP_H / 2
app_right = APP_L + APP_W

entra_cx = ENTRA_L + ENTRA_W / 2
entra_bottom = ENTRA_T + ENTRA_H
entra_left = ENTRA_L

# ── ① Client → Entra ID (dashed yellow, public) ──
arrow(sz_cx, SZ_T, entra_cx + Inches(0.3), entra_bottom, YELLOW, Pt(2), dashed=True)
arrow(dl_cx, DL_T, entra_cx - Inches(0.3), entra_bottom, YELLOW, Pt(2), dashed=True)
step(Inches(0.2), Inches(0.62), "①", "Entra ID Auth\n⚠️ Public", YELLOW)

# ── ② Entra ID → Client (dashed yellow, public) ──
arrow(entra_left, entra_bottom, sz_cx - Inches(0.2), SZ_T, YELLOW, Pt(2), dashed=True)
arrow(entra_left - Inches(0.1), entra_bottom, dl_cx - Inches(0.2), DL_T, YELLOW, Pt(2), dashed=True)
step(Inches(2.2), Inches(0.62), "②", "Token Returned\n⚠️ Public", YELLOW)

# ── ③ Client → SD-WAN (blue) ──
arrow(op_right, sz_cy, sdwan_left, SDWAN_CY - Inches(0.15), STEP_BLUE)
arrow(op_right, dl_cy, sdwan_left, SDWAN_CY + Inches(0.15), STEP_BLUE)
step(Inches(2.15), sz_cy - Inches(0.25), "③", "Token · SD-WAN", STEP_BLUE)

# ── ③ SD-WAN → VNet Gateway (green) ──
arrow(sdwan_right, SDWAN_CY, gw_left, gw_cy, SDWAN_GREEN, Pt(3))
step(Inches(3.3), Inches(2.5), "③", "SD-WAN Transit", STEP_BLUE)

# ── ④ VNet Gateway → PE Workspace ──
arrow(gw_cx - Inches(0.3), gw_bottom, pews_cx, pews_top, STEP_BLUE)
step(CE2_L + Inches(0.05), Inches(1.45), "④", "Feed Discovery", STEP_BLUE)

# ── ⑤ PE Workspace → PE Host Pool ──
arrow(PEWS_L + PEWS_W, pews_cy, PEHP_L, pehp_cy, STEP_BLUE)
step(PEWS_L + PEWS_W - Inches(0.1), PEWS_T + PEWS_H + Inches(0.02), "⑤", "Session Request", STEP_BLUE)

# ── ⑥ PE Host Pool → Session Host (Reverse Connect) ──
arrow(pehp_right, pehp_cy, sh_left, sh_cy, STEP_BLUE)
step(CN3_L - Inches(0.6), pehp_cy - Inches(0.25), "⑥", "Reverse Connect", STEP_BLUE)

# ── ⑦ Session Host → App Server (VNet Peering) ──
arrow(sh_left, sh_cy + Inches(0.15), app_right, app_cy - Inches(0.1), STEP_BLUE)
step(CN3_L - Inches(0.55), APP_T - Inches(0.15), "⑦", "VNet Peering", STEP_BLUE)

# ── ⑧ App Server → Session Host ──
arrow(app_right, app_cy + Inches(0.1), sh_left, sh_cy + Inches(0.35), STEP_BLUE)
step(CN3_L - Inches(0.55), APP_T + Inches(0.35), "⑧", "Return Data", STEP_BLUE)

# ── ⑨ RDP Shortpath: Session Host → GW ──
arrow(sh_left, sh_top + Inches(0.1), gw_right, gw_cy, SHORTPATH_ORANGE, Pt(3))
step(CN3_L - Inches(0.6), GW_T - Inches(0.15), "⑨", "Shortpath · UDP", STEP_RED)

# ── ⑨ GW → SD-WAN ──
arrow(gw_left, gw_cy - Inches(0.05), sdwan_right, SDWAN_CY - Inches(0.2), SHORTPATH_ORANGE, Pt(3))
step(Inches(3.3), GW_T - Inches(0.15), "⑨", "SD-WAN Return", STEP_RED)

# ── ⑨ SD-WAN → Clients ──
arrow(sdwan_left, SDWAN_CY - Inches(0.3), op_right, sz_cy + Inches(0.15), SHORTPATH_ORANGE, Pt(3))
arrow(sdwan_left, SDWAN_CY + Inches(0.3), op_right, dl_cy + Inches(0.15), SHORTPATH_ORANGE, Pt(3))

# ══════════════════════════════════════════════════════════
# LEGEND (bottom-left)
# ══════════════════════════════════════════════════════════
LG_L = Inches(0.2)
LG_T = Inches(6.55)
txt(LG_L, LG_T, Inches(2), Inches(0.2), "Legend:", fs=8, fc=BLACK, bold=True, align=PP_ALIGN.LEFT)

y = LG_T + Inches(0.25)
slide.shapes.add_connector(1, LG_L, y, LG_L + Inches(0.5), y).line.color.rgb = YELLOW
txt(LG_L + Inches(0.55), y - Inches(0.08), Inches(2.5), Inches(0.2),
    "①② Entra ID Auth (⚠️ Public Internet)", fs=7)

y += Inches(0.22)
c = slide.shapes.add_connector(1, LG_L, y, LG_L + Inches(0.5), y)
c.line.color.rgb = STEP_BLUE; c.line.width = Pt(2)
txt(LG_L + Inches(0.55), y - Inches(0.08), Inches(2.5), Inches(0.2),
    "③~⑧ Private Connection Flow", fs=7)

y += Inches(0.22)
c = slide.shapes.add_connector(1, LG_L, y, LG_L + Inches(0.5), y)
c.line.color.rgb = SHORTPATH_ORANGE; c.line.width = Pt(3)
txt(LG_L + Inches(0.55), y - Inches(0.08), Inches(2.5), Inches(0.2),
    "⑨ RDP Shortpath (UDP · Private)", fs=7)

y += Inches(0.22)
c = slide.shapes.add_connector(1, LG_L, y, LG_L + Inches(0.5), y)
c.line.color.rgb = SDWAN_GREEN; c.line.width = Pt(3)
txt(LG_L + Inches(0.55), y - Inches(0.08), Inches(2.5), Inches(0.2),
    "SD-WAN Private Transit", fs=7)

# ══════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════
out = "/Users/tema/projects/avd-privatelink/AVD-PrivateLink-Architecture-v3.pptx"
prs.save(out)
print(f"✅ PPT v3 saved to: {out}")
