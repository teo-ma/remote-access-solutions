#!/usr/bin/env python3
"""Generate an editable PowerPoint slide for the AVD Private Link architecture."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

# ── Color palette ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_BLUE = RGBColor(0x0D, 0x47, 0xA1)
BLUE = RGBColor(0x15, 0x65, 0xC0)
LIGHT_BLUE = RGBColor(0x1E, 0x88, 0xE5)
OCEAN_BLUE = RGBColor(0x02, 0x77, 0xBD)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
DARK_ORANGE = RGBColor(0xBF, 0x36, 0x0C)
LIGHT_ORANGE = RGBColor(0xEF, 0x6C, 0x00)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
DARK_GREEN = RGBColor(0x1B, 0x5E, 0x20)
YELLOW = RGBColor(0xF9, 0xA8, 0x25)
AMBER_STROKE = RGBColor(0xE6, 0x51, 0x00)
STEP_BLUE = RGBColor(0x1A, 0x73, 0xE8)
STEP_RED = RGBColor(0xD5, 0x00, 0x00)
SHORTPATH_ORANGE = RGBColor(0xFF, 0x6D, 0x00)
SDWAN_GREEN = RGBColor(0x00, 0xC8, 0x53)


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color,
                     text="", font_size=10, font_color=WHITE, bold=False,
                     border_width=Pt(2), transparency=0):
    """Add a rounded rectangle shape with text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = border_width

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=9,
                 font_color=BLACK, bold=False, alignment=PP_ALIGN.CENTER):
    """Add an editable text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    return txBox


def add_oval(slide, left, top, width, height, fill_color, border_color,
             text="", font_size=9, font_color=WHITE):
    """Add an oval/circle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = True
    return shape


def add_connector(slide, start_x, start_y, end_x, end_y, color, width=Pt(2)):
    """Add a straight connector line (as freeform with arrow)."""
    connector = slide.shapes.add_connector(
        1,  # msoConnectorStraight
        start_x, start_y, end_x, end_y
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def add_arrow(slide, start_x, start_y, end_x, end_y, color, width=Pt(2)):
    """Add a connector with arrowhead."""
    conn = add_connector(slide, start_x, start_y, end_x, end_y, color, width)
    # Add end arrow
    conn.line.fill.solid()
    conn.line.color.rgb = color
    end_arrow = conn.line._get_or_add_ln()
    from lxml import etree
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    tail_end = etree.SubElement(end_arrow, '{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
    tail_end.set('type', 'triangle')
    tail_end.set('w', 'med')
    tail_end.set('len', 'med')
    return conn


def add_step_label(slide, x, y, num, text, color=STEP_BLUE):
    """Add a circled step number with description text."""
    # Number circle
    circle = add_oval(slide, x, y, Inches(0.32), Inches(0.32),
                      color, color, num, font_size=9, font_color=WHITE)
    # Label text
    label = add_text_box(slide, x + Inches(0.36), y - Inches(0.02),
                         Inches(1.6), Inches(0.36), text,
                         font_size=8, font_color=BLACK, bold=False)
    return circle, label


# ══════════════════════════════════════════════════════════════════
# Title
# ══════════════════════════════════════════════════════════════════
add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12), Inches(0.5),
             "Azure Virtual Desktop — Private Link Architecture",
             font_size=22, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.LEFT)
add_text_box(slide, Inches(0.3), Inches(0.55), Inches(10), Inches(0.3),
             "End-to-end private connectivity via SD-WAN + AVD Private Link  |  Azure China (21Vianet)",
             font_size=10, font_color=RGBColor(0x66, 0x66, 0x66), bold=False, alignment=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════
# Region boxes (background)
# ══════════════════════════════════════════════════════════════════

# On-Premises Offices (left side)
onprem_l, onprem_t = Inches(0.3), Inches(1.2)
onprem_w, onprem_h = Inches(2.2), Inches(4.0)
add_rounded_rect(slide, onprem_l, onprem_t, onprem_w, onprem_h,
                 GREEN, DARK_GREEN, "", border_width=Pt(3))
add_text_box(slide, onprem_l, onprem_t + Inches(0.1), onprem_w, Inches(0.35),
             "🏢 On-Premises Offices", font_size=11, font_color=WHITE, bold=True)

# Thin Client - Shenzhen
add_rounded_rect(slide, onprem_l + Inches(0.2), onprem_t + Inches(0.6),
                 Inches(1.8), Inches(1.3), RGBColor(0x38, 0x8E, 0x3C),
                 DARK_GREEN, "🖥️ Thin Clients\n(Shenzhen)\n🔒 10.1.1.0/24",
                 font_size=9, font_color=WHITE, bold=False, border_width=Pt(1))

# Thin Client - Dalian
add_rounded_rect(slide, onprem_l + Inches(0.2), onprem_t + Inches(2.2),
                 Inches(1.8), Inches(1.3), RGBColor(0x38, 0x8E, 0x3C),
                 DARK_GREEN, "🖥️ Thin Clients\n(Dalian)\n🔒 10.2.1.0/24",
                 font_size=9, font_color=WHITE, bold=False, border_width=Pt(1))

# SD-WAN circle
sdwan_cx, sdwan_cy = Inches(3.4), Inches(2.8)
sdwan_r = Inches(0.7)
add_oval(slide, sdwan_cx - sdwan_r, sdwan_cy - sdwan_r,
         sdwan_r * 2, sdwan_r * 2, RGBColor(0x00, 0x60, 0x00),
         DARK_GREEN, "🔒\nSD-WAN\nPrivate\nNetwork",
         font_size=8, font_color=WHITE)

# Azure China East 2
ce2_l, ce2_t = Inches(5.0), Inches(1.0)
ce2_w, ce2_h = Inches(4.0), Inches(5.8)
add_rounded_rect(slide, ce2_l, ce2_t, ce2_w, ce2_h,
                 BLUE, DARK_BLUE, "", border_width=Pt(3))
add_text_box(slide, ce2_l, ce2_t + Inches(0.08), ce2_w, Inches(0.35),
             "☁️ Azure China East 2", font_size=12, font_color=WHITE, bold=True)

# VNet Gateway (top of CE2)
gw_l = ce2_l + Inches(0.5)
gw_t = ce2_t + Inches(0.5)
gw_w = Inches(3.0)
gw_h = Inches(0.6)
add_rounded_rect(slide, gw_l, gw_t, gw_w, gw_h,
                 DARK_BLUE, RGBColor(0x01, 0x57, 0x9B),
                 "🔗 VNet Gateway · SD-WAN Hub",
                 font_size=10, font_color=WHITE, bold=True, border_width=Pt(2))

# AVD Control Plane sub-box
avd_l = ce2_l + Inches(0.3)
avd_t = ce2_t + Inches(1.4)
avd_w = Inches(3.4)
avd_h = Inches(2.0)
add_rounded_rect(slide, avd_l, avd_t, avd_w, avd_h,
                 LIGHT_BLUE, BLUE, "", border_width=Pt(1))
add_text_box(slide, avd_l, avd_t + Inches(0.05), avd_w, Inches(0.3),
             "AVD Control Plane · Private Link", font_size=9, font_color=WHITE, bold=True)

# PE Workspace
pe_ws_l = avd_l + Inches(0.15)
pe_ws_t = avd_t + Inches(0.45)
pe_ws_w, pe_ws_h = Inches(1.45), Inches(1.3)
add_rounded_rect(slide, pe_ws_l, pe_ws_t, pe_ws_w, pe_ws_h,
                 RGBColor(0x42, 0xA5, 0xF5), LIGHT_BLUE,
                 "🔗\nPrivate Endpoint\nWorkspace\nFeed Discovery",
                 font_size=8, font_color=WHITE, border_width=Pt(1))

# PE Host Pool
pe_hp_l = avd_l + Inches(1.75)
pe_hp_t = avd_t + Inches(0.45)
pe_hp_w, pe_hp_h = Inches(1.45), Inches(1.3)
add_rounded_rect(slide, pe_hp_l, pe_hp_t, pe_hp_w, pe_hp_h,
                 RGBColor(0x42, 0xA5, 0xF5), LIGHT_BLUE,
                 "🔗\nPrivate Endpoint\nHost Pool\nRDP Connection",
                 font_size=8, font_color=WHITE, border_width=Pt(1))

# Business App VNet sub-box
biz_l = ce2_l + Inches(0.3)
biz_t = ce2_t + Inches(3.6)
biz_w = Inches(3.4)
biz_h = Inches(1.8)
add_rounded_rect(slide, biz_l, biz_t, biz_w, biz_h,
                 OCEAN_BLUE, RGBColor(0x01, 0x57, 0x9B), "", border_width=Pt(1))
add_text_box(slide, biz_l, biz_t + Inches(0.05), biz_w, Inches(0.3),
             "Business App VNet · 10.10.0.0/16", font_size=9, font_color=WHITE, bold=True)

# App Server
add_rounded_rect(slide, biz_l + Inches(0.4), biz_t + Inches(0.5),
                 Inches(2.6), Inches(1.0),
                 RGBColor(0x03, 0x9B, 0xE5), OCEAN_BLUE,
                 "🖧 Business Application Server",
                 font_size=10, font_color=WHITE, bold=False, border_width=Pt(1))

# Azure China North 3
cn3_l, cn3_t = Inches(9.5), Inches(1.0)
cn3_w, cn3_h = Inches(3.5), Inches(5.8)
add_rounded_rect(slide, cn3_l, cn3_t, cn3_w, cn3_h,
                 ORANGE, DARK_ORANGE, "", border_width=Pt(3))
add_text_box(slide, cn3_l, cn3_t + Inches(0.08), cn3_w, Inches(0.35),
             "☁️ Azure China North 3", font_size=12, font_color=WHITE, bold=True)

# Session Host VNet sub-box
sh_vnet_l = cn3_l + Inches(0.2)
sh_vnet_t = cn3_t + Inches(0.5)
sh_vnet_w = Inches(3.1)
sh_vnet_h = Inches(5.0)
add_rounded_rect(slide, sh_vnet_l, sh_vnet_t, sh_vnet_w, sh_vnet_h,
                 LIGHT_ORANGE, ORANGE, "", border_width=Pt(1))
add_text_box(slide, sh_vnet_l, sh_vnet_t + Inches(0.05), sh_vnet_w, Inches(0.3),
             "Session Host VNet · 10.20.0.0/16", font_size=9, font_color=WHITE, bold=True)

# Session Host VM
add_rounded_rect(slide, sh_vnet_l + Inches(0.2), sh_vnet_t + Inches(0.5),
                 Inches(2.7), Inches(1.2),
                 RGBColor(0xF4, 0x7B, 0x13), ORANGE,
                 "🖥️ Windows Session Host VM\nDesktop & RemoteApp",
                 font_size=10, font_color=WHITE, bold=True, border_width=Pt(1))

# VNet Peering
add_rounded_rect(slide, sh_vnet_l + Inches(0.2), sh_vnet_t + Inches(1.9),
                 Inches(2.7), Inches(0.8),
                 LIGHT_ORANGE, ORANGE,
                 "🔗 VNet Peering ↔ China East 2",
                 font_size=9, font_color=WHITE, bold=False, border_width=Pt(1))

# Private DNS Zone
add_rounded_rect(slide, sh_vnet_l + Inches(0.2), sh_vnet_t + Inches(2.9),
                 Inches(2.7), Inches(0.8),
                 DARK_ORANGE, ORANGE,
                 "🔒 Private DNS Zone\nprivatelink.wvd.azure.cn",
                 font_size=9, font_color=WHITE, bold=False, border_width=Pt(1))

# NSG
add_rounded_rect(slide, sh_vnet_l + Inches(0.2), sh_vnet_t + Inches(3.9),
                 Inches(2.7), Inches(0.8),
                 DARK_ORANGE, ORANGE,
                 "🛡️ NSG · Outbound Rules\nAllow UDP 3390",
                 font_size=9, font_color=WHITE, bold=False, border_width=Pt(1))

# Azure Entra ID (top center)
entra_l = Inches(5.5)
entra_t = Inches(0.1)
entra_w, entra_h = Inches(2.5), Inches(0.65)
add_rounded_rect(slide, entra_l, entra_t, entra_w, entra_h,
                 YELLOW, AMBER_STROKE,
                 "🌐 Azure Entra ID (Public Endpoint)",
                 font_size=9, font_color=BLACK, bold=True, border_width=Pt(2))

# ══════════════════════════════════════════════════════════════════
# Connection arrows & step labels
# ══════════════════════════════════════════════════════════════════

# Define key anchor points (centers)
sz_r = onprem_l + onprem_w      # right edge of Shenzhen
dl_r = onprem_l + onprem_w      # right edge of Dalian
sz_cy = onprem_t + Inches(1.25)
dl_cy = onprem_t + Inches(2.85)
sdwan_l_edge = sdwan_cx - sdwan_r
sdwan_r_edge = sdwan_cx + sdwan_r
gw_cx = gw_l + gw_w / 2
gw_cy_mid = gw_t + gw_h / 2
pe_ws_cx = pe_ws_l + pe_ws_w / 2
pe_ws_cy = pe_ws_t + pe_ws_h / 2
pe_hp_cx = pe_hp_l + pe_hp_w / 2
pe_hp_cy = pe_hp_t + pe_hp_h / 2
sh_cx = sh_vnet_l + Inches(1.55)
sh_cy_top = sh_vnet_t + Inches(0.5)
sh_cy_bot = sh_vnet_t + Inches(1.7)
app_cx = biz_l + Inches(1.7)
app_cy = biz_t + Inches(1.0)

# ① Shenzhen → SD-WAN
add_arrow(slide, sz_r, sz_cy, sdwan_l_edge, sdwan_cy - Inches(0.2), STEP_BLUE)
add_step_label(slide, Inches(2.55), sz_cy - Inches(0.35), "①", "Initiate Connection", STEP_BLUE)

# ① Dalian → SD-WAN
add_arrow(slide, dl_r, dl_cy, sdwan_l_edge, sdwan_cy + Inches(0.2), STEP_BLUE)
add_step_label(slide, Inches(2.55), dl_cy - Inches(0.0), "①", "Initiate Connection", STEP_BLUE)

# ② SD-WAN → VNet Gateway
add_arrow(slide, sdwan_r_edge, sdwan_cy, gw_l, gw_cy_mid, SDWAN_GREEN, Pt(3))
add_step_label(slide, Inches(4.3), Inches(1.9), "②", "SD-WAN Private Transit", STEP_BLUE)

# ③ VNet Gateway → Entra ID
add_arrow(slide, gw_cx, gw_t, entra_l + entra_w / 2, entra_t + entra_h, STEP_BLUE)
add_step_label(slide, Inches(6.0), Inches(0.82), "③", "Entra ID Auth", STEP_BLUE)

# ④ Entra ID → PE Workspace
add_arrow(slide, entra_l + entra_w, entra_t + entra_h / 2,
          pe_ws_cx, pe_ws_t, STEP_BLUE)
add_step_label(slide, Inches(7.1), Inches(0.15), "④", "Auth Token Issued", STEP_BLUE)

# ⑤ PE Workspace → PE Host Pool
add_arrow(slide, pe_ws_l + pe_ws_w, pe_ws_cy, pe_hp_l, pe_hp_cy, STEP_BLUE)
add_step_label(slide, Inches(6.6), Inches(2.8), "⑤", "Feed Discovery", STEP_BLUE)

# ⑥ PE Host Pool → Session Host
add_arrow(slide, pe_hp_l + pe_hp_w, pe_hp_cy, sh_vnet_l, sh_cy_top + Inches(0.5), STEP_BLUE)
add_step_label(slide, Inches(8.7), Inches(2.0), "⑥", "RDP Reverse Connect", STEP_BLUE)

# ⑦ Session Host → App Server
add_arrow(slide, sh_vnet_l, sh_cy_bot, biz_l + biz_w, app_cy - Inches(0.15), STEP_BLUE)
add_step_label(slide, Inches(8.3), Inches(4.0), "⑦", "VNet Peering to App", STEP_BLUE)

# ⑧ App Server → Session Host
add_arrow(slide, biz_l + biz_w, app_cy + Inches(0.15), sh_vnet_l, sh_cy_bot + Inches(0.3), STEP_BLUE)
add_step_label(slide, Inches(8.3), Inches(4.55), "⑧", "Return Data", STEP_BLUE)

# ⑨ Session Host → VNet Gateway (Shortpath)
add_arrow(slide, sh_vnet_l, sh_cy_top + Inches(0.2),
          gw_l + gw_w, gw_cy_mid, SHORTPATH_ORANGE, Pt(3))
add_step_label(slide, Inches(8.3), Inches(1.15), "⑨", "RDP Shortpath · UDP", STEP_RED)

# ⑨ VNet Gateway → SD-WAN
add_arrow(slide, gw_l, gw_cy_mid, sdwan_r_edge, sdwan_cy - Inches(0.3), SHORTPATH_ORANGE, Pt(3))
add_step_label(slide, Inches(3.9), Inches(1.2), "⑨", "SD-WAN Transit", STEP_RED)

# ⑨ SD-WAN → Shenzhen
add_arrow(slide, sdwan_l_edge, sdwan_cy - Inches(0.35), sz_r, sz_cy + Inches(0.2), SHORTPATH_ORANGE, Pt(3))

# ⑨ SD-WAN → Dalian
add_arrow(slide, sdwan_l_edge, sdwan_cy + Inches(0.35), dl_r, dl_cy + Inches(0.2), SHORTPATH_ORANGE, Pt(3))

# ══════════════════════════════════════════════════════════════════
# Legend
# ══════════════════════════════════════════════════════════════════
legend_l = Inches(0.3)
legend_t = Inches(5.6)
add_text_box(slide, legend_l, legend_t, Inches(3), Inches(0.3),
             "Legend:", font_size=9, font_color=BLACK, bold=True, alignment=PP_ALIGN.LEFT)

# Blue arrow = normal flow
add_connector(slide, legend_l + Inches(0.05), legend_t + Inches(0.4),
              legend_l + Inches(0.6), legend_t + Inches(0.4), STEP_BLUE, Pt(2))
add_text_box(slide, legend_l + Inches(0.65), legend_t + Inches(0.28),
             Inches(2), Inches(0.25),
             "①~⑧ Initial Connection Flow", font_size=8, font_color=BLACK)

# Orange arrow = shortpath
add_connector(slide, legend_l + Inches(0.05), legend_t + Inches(0.7),
              legend_l + Inches(0.6), legend_t + Inches(0.7), SHORTPATH_ORANGE, Pt(3))
add_text_box(slide, legend_l + Inches(0.65), legend_t + Inches(0.58),
             Inches(2), Inches(0.25),
             "⑨ RDP Shortpath (Direct UDP)", font_size=8, font_color=BLACK)

# Green arrow = SD-WAN
add_connector(slide, legend_l + Inches(0.05), legend_t + Inches(1.0),
              legend_l + Inches(0.6), legend_t + Inches(1.0), SDWAN_GREEN, Pt(3))
add_text_box(slide, legend_l + Inches(0.65), legend_t + Inches(0.88),
             Inches(2), Inches(0.25),
             "SD-WAN Private Transit (Inbound)", font_size=8, font_color=BLACK)

# Yellow = Public
add_rounded_rect(slide, legend_l, legend_t + Inches(1.2),
                 Inches(0.6), Inches(0.25), YELLOW, AMBER_STROKE, "",
                 border_width=Pt(1))
add_text_box(slide, legend_l + Inches(0.65), legend_t + Inches(1.18),
             Inches(2), Inches(0.25),
             "⚠️ Public Endpoint (Entra ID only)", font_size=8, font_color=BLACK)

# ══════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════
output_path = "/Users/tema/projects/avd-privatelink/AVD-PrivateLink-Architecture.pptx"
prs.save(output_path)
print(f"✅ PPT saved to: {output_path}")
